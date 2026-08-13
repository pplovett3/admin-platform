# -*- coding: utf-8 -*-
"""
昌吉州AI智能体大赛 路演PPT生成脚本（美化版 v2）
项目：AI数字人工匠导师智能体 - 职教元宇宙解决方案
美术风格：深空科技风 + 玻璃拟态 + 几何装饰 + 数据可视化
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import math

# ============ 设计令牌 v2 ============
# 主色系：深空蓝 + 科技青 + 活力橙 + 紫罗兰
C_INK       = RGBColor(0x0A, 0x14, 0x28)   # 墨黑蓝（最深背景）
C_NAVY      = RGBColor(0x10, 0x2A, 0x4A)   # 海军蓝
C_ROYAL     = RGBColor(0x1E, 0x5C, 0xE5)   # 皇家蓝
C_CYAN      = RGBColor(0x00, 0xD4, 0xFF)   # 电光青
C_TEAL      = RGBColor(0x14, 0xB8, 0xA6)   # 青绿
C_ORANGE    = RGBColor(0xFF, 0x8A, 0x3D)   # 暖橙
C_AMBER     = RGBColor(0xFF, 0xC5, 0x3D)   # 琥珀
C_VIOLET    = RGBColor(0x8B, 0x5C, 0xF6)   # 紫罗兰
C_MAGENTA   = RGBColor(0xEC, 0x48, 0x99)   # 品红
C_PAPER     = RGBColor(0xF4, 0xF6, 0xFA)   # 纸白
C_CARD      = RGBColor(0xFF, 0xFF, 0xFF)   # 卡片白
C_CARD_DARK = RGBColor(0x14, 0x24, 0x3F)   # 深卡片
C_CARD_MID  = RGBColor(0x1B, 0x2E, 0x4D)   # 中卡片
C_TEXT      = RGBColor(0x1A, 0x22, 0x33)   # 主文本
C_TEXT_LT   = RGBColor(0xFF, 0xFF, 0xFF)   # 反白
C_TEXT_SUB  = RGBColor(0x6B, 0x7A, 0x8F)   # 次要文本
C_TEXT_DIM  = RGBColor(0x9A, 0xA8, 0xBC)   # 暗淡文本
C_LINE      = RGBColor(0xE2, 0xE8, 0xF0)   # 浅分割线
C_LINE_DARK = RGBColor(0x2A, 0x3D, 0x5F)   # 深分割线
C_GRID      = RGBColor(0x1C, 0x2D, 0x4D)   # 网格线（深色背景）

FONT_T = "微软雅黑"
FONT_B = "微软雅黑"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# 配色组：用于卡片循环
PALETTE = [C_CYAN, C_ORANGE, C_ROYAL, C_VIOLET, C_TEAL, C_MAGENTA]


# ============ 底层工具 ============
def _set_ea_font(run, font=FONT_B):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font)
    latin = rPr.find(qn('a:latin'))
    if latin is None:
        latin = etree.SubElement(rPr, qn('a:latin'))
    latin.set('typeface', font)

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def set_gradient(shape, c1, c2, angle=45, c1_pos=0, c2_pos=100):
    """三色渐变可选；这里用两色"""
    sp = shape.fill._xPr
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    grad = etree.SubElement(sp, qn('a:gradFill'))
    grad.set('flip', 'none'); grad.set('rotWithShape', '1')
    lst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, col in [(c1_pos, c1), (c2_pos, c2)]:
        gs = etree.SubElement(lst, qn('a:gs'))
        gs.set('pos', str(int(pos * 1000)))
        c = etree.SubElement(gs, qn('a:srgbClr'))
        c.set('val', '%02X%02X%02X' % (col[0], col[1], col[2]))
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(int(angle * 60000)))
    lin.set('scaled', '1')

def set_gradient3(shape, c1, c2, c3, angle=45):
    """三色渐变"""
    sp = shape.fill._xPr
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    grad = etree.SubElement(sp, qn('a:gradFill'))
    grad.set('flip', 'none'); grad.set('rotWithShape', '1')
    lst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, col in [(0, c1), (50, c2), (100, c3)]:
        gs = etree.SubElement(lst, qn('a:gs'))
        gs.set('pos', str(int(pos * 1000)))
        c = etree.SubElement(gs, qn('a:srgbClr'))
        c.set('val', '%02X%02X%02X' % (col[0], col[1], col[2]))
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(int(angle * 60000)))
    lin.set('scaled', '1')

def set_shadow(shape, blur=Pt(8), dist=Pt(3), dir=5400000, alpha=70):
    """柔和阴影：dir单位1/60000度，5400000=90度（向下）"""
    spPr = shape._element.spPr
    # 移除已有 effectLst
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    effect = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(effect, qn('a:outerShdw'))
    sh.set('blurRad', str(int(blur)))
    sh.set('dist', str(int(dist)))
    sh.set('dir', str(int(dir)))
    sh.set('rotWithShape', '0')
    clr = etree.SubElement(sh, qn('a:srgbClr'))
    clr.set('val', '0A1428')
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(int((100 - alpha) * 1000)))

def set_glow(shape, color, rad=Pt(12), alpha=40):
    """外发光"""
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    effect = etree.SubElement(spPr, qn('a:effectLst'))
    gl = etree.SubElement(effect, qn('a:glow'))
    gl.set('rad', str(int(rad)))
    clr = etree.SubElement(gl, qn('a:srgbClr'))
    clr.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(int(alpha * 1000)))

def no_line(shape):
    shape.line.fill.background()

def set_line(shape, color, w=Pt(1)):
    shape.line.color.rgb = color
    shape.line.width = w

def set_round_radius(shape, ratio=0.06):
    try:
        shape.adjustments[0] = ratio
    except Exception:
        pass

def add_rect(s, x, y, w, h, color=None, gradient=None, grad_angle=45, gradient3=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if gradient3:
        set_gradient3(shp, gradient3[0], gradient3[1], gradient3[2], grad_angle)
        no_line(shp)
    elif gradient:
        set_gradient(shp, gradient[0], gradient[1], grad_angle)
        no_line(shp)
    elif color is not None:
        set_fill(shp, color)
    else:
        shp.fill.background(); no_line(shp)
    return shp

def add_rrect(s, x, y, w, h, color=None, line_color=None, line_w=Pt(1), radius=0.06,
              gradient=None, grad_angle=45, shadow=False, glow=None):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    set_round_radius(shp, radius)
    if gradient:
        set_gradient(shp, gradient[0], gradient[1], grad_angle); no_line(shp)
    elif color is not None:
        set_fill(shp, color)
    else:
        shp.fill.background()
    if line_color is not None:
        set_line(shp, line_color, line_w)
    else:
        no_line(shp)
    if shadow:
        set_shadow(shp)
    if glow:
        set_glow(shp, glow[0], rad=glow[1], alpha=glow[2])
    return shp

def add_oval(s, x, y, w, h, color=None, line_color=None, line_w=Pt(1),
             gradient=None, grad_angle=45, glow=None, shadow=False):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.shadow.inherit = False
    if gradient:
        set_gradient(shp, gradient[0], gradient[1], grad_angle); no_line(shp)
    elif color is not None:
        set_fill(shp, color)
    else:
        shp.fill.background()
    if line_color is not None:
        set_line(shp, line_color, line_w)
    else:
        no_line(shp)
    if glow:
        set_glow(shp, glow[0], rad=glow[1], alpha=glow[2])
    if shadow:
        set_shadow(shp)
    return shp

def add_text(s, x, y, w, h, text, size=14, color=C_TEXT, bold=False,
             font=FONT_B, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.2, spacing=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if spacing:
            p.space_after = Pt(spacing)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
        _set_ea_font(r, font)
    return tb

def add_bullets(s, x, y, w, h, items, size=13, color=C_TEXT,
                bullet_color=C_CYAN, line_spacing=1.5, font=FONT_B,
                bullet_char="●", gap=6):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(gap)
        r1 = p.add_run()
        r1.text = bullet_char + "  "
        r1.font.size = Pt(size)
        r1.font.name = font
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        _set_ea_font(r1, font)
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.name = font
        r2.font.color.rgb = color
        _set_ea_font(r2, font)
    return tb


# ============ 装饰组件 ============
def deco_grid_bg(s, color=C_GRID, spacing=Inches(0.5), alpha=15):
    """深色页面的点阵网格背景"""
    cols = int(13.333 / 0.5) + 2
    rows = int(7.5 / 0.5) + 2
    for r in range(rows):
        for c in range(cols):
            dot = add_oval(s, Inches(c * 0.5), Inches(r * 0.5),
                           Pt(1.5), Pt(1.5), color=color)
            # 设置透明度
            spPr = dot._element.spPr
            sf = spPr.find(qn('a:solidFill'))
            if sf is not None:
                clr = sf.find(qn('a:srgbClr'))
                if clr is not None:
                    a = etree.SubElement(clr, qn('a:alpha'))
                    a.set('val', str(int(alpha * 1000)))

def deco_diagonal_lines(s, x, y, w, h, color, count=5, alpha=20):
    """对角装饰线"""
    for i in range(count):
        offset = Inches(i * 0.3)
        ln = add_rect(s, x + offset, y, Pt(1), h, color=color)
        spPr = ln._element.spPr
        sf = spPr.find(qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                a = etree.SubElement(clr, qn('a:alpha'))
                a.set('val', str(int(alpha * 1000)))

def deco_corner_rings(s, x, y, size, color, line_w=Pt(1.5), alpha=60):
    """装饰圆环"""
    ring = add_oval(s, x, y, size, size, line_color=color, line_w=line_w)
    spPr = ring._element.spPr
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        clr = ln.find(qn('a:srgbClr'))
        if clr is not None:
            a = etree.SubElement(clr, qn('a:alpha'))
            a.set('val', str(int(alpha * 1000)))
    return ring

def deco_glow_dot(s, x, y, size, color, alpha=50):
    """光晕点"""
    dot = add_oval(s, x, y, size, size, color=color)
    set_glow(dot, color, rad=Pt(18), alpha=alpha)
    return dot

def set_alpha(shape, alpha_pct):
    """给纯色填充形状设置透明度（0-100，100=完全不透明）"""
    spPr = shape._element.spPr
    sf = spPr.find(qn('a:solidFill'))
    if sf is not None:
        clr = sf.find(qn('a:srgbClr'))
        if clr is not None:
            a = etree.SubElement(clr, qn('a:alpha'))
            a.set('val', str(int(alpha_pct * 1000)))


# ============ 通用页头/页脚 ============
def page_header(s, kicker, title, num=None, total=19, dark=False):
    """内容页统一页头"""
    # 顶部细装饰条
    add_rect(s, 0, 0, SLIDE_W, Pt(4))
    set_gradient(add_rect(s, 0, 0, SLIDE_W, Pt(4)),
                 (0x00, 0xD4, 0xFF), (0xFF, 0x8A, 0x3D), angle=0)
    # 左侧色条标记
    bar = add_rect(s, Inches(0.6), Inches(0.6), Inches(0.06), Inches(0.55), color=C_CYAN)
    # 小标签
    add_text(s, Inches(0.82), Inches(0.55), Inches(9), Inches(0.3),
             kicker, size=11, color=C_CYAN, bold=True)
    # 大标题
    title_color = C_TEXT_LT if dark else C_TEXT
    add_text(s, Inches(0.82), Inches(0.82), Inches(11.5), Inches(0.55),
             title, size=26, color=title_color, bold=True)
    # 标题下短装饰线（双色）
    add_rect(s, Inches(0.82), Inches(1.42), Inches(0.5), Pt(3), color=C_ORANGE)
    add_rect(s, Inches(1.32), Inches(1.42), Inches(0.25), Pt(3), color=C_CYAN)
    if num is not None:
        add_text(s, Inches(12.2), Inches(7.05), Inches(1.0), Inches(0.3),
                 f"{num:02d} / {total:02d}", size=10,
                 color=C_TEXT_DIM if dark else C_TEXT_SUB, align=PP_ALIGN.RIGHT)

def page_footer(s, dark=False):
    col = C_LINE_DARK if dark else C_LINE
    add_rect(s, Inches(0.6), Inches(7.15), Inches(12.1), Pt(0.5), color=col)
    add_text(s, Inches(0.6), Inches(7.22), Inches(6), Inches(0.25),
             "AI 数字人工匠导师智能体  ·  昌吉州 AI 智能体大赛",
             size=9, color=C_TEXT_DIM if dark else C_TEXT_SUB)


# ============ 1. 封面 ============
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    # 深色渐变背景（三色）
    bg = add_rect(s, 0, 0, SLIDE_W, SLIDE_H)
    set_gradient3(bg, (0x05, 0x0A, 0x1A), (0x0A, 0x1F, 0x3D), (0x12, 0x2E, 0x5C), angle=120)
    # 点阵网格
    deco_grid_bg(s, color=RGBColor(0x1A, 0x2E, 0x4D), alpha=12)
    # 右侧大装饰圆环组
    deco_corner_rings(s, Inches(10.2), Inches(0.8), Inches(2.8), C_CYAN, line_w=Pt(1.5), alpha=40)
    deco_corner_rings(s, Inches(10.7), Inches(1.3), Inches(1.8), C_ORANGE, line_w=Pt(1), alpha=50)
    deco_corner_rings(s, Inches(11.0), Inches(1.6), Inches(1.2), C_VIOLET, line_w=Pt(0.75), alpha=40)
    # 光晕点
    deco_glow_dot(s, Inches(11.5), Inches(2.2), Inches(0.18), C_CYAN, alpha=60)
    deco_glow_dot(s, Inches(12.3), Inches(4.5), Inches(0.14), C_ORANGE, alpha=55)
    deco_glow_dot(s, Inches(1.5), Inches(6.5), Inches(0.12), C_VIOLET, alpha=50)
    # 左下装饰对角线
    deco_diagonal_lines(s, Inches(0.5), Inches(5.8), Inches(4), Inches(1.5),
                        C_CYAN, count=8, alpha=15)
    # 顶部赛事胶囊标签
    tag = add_rrect(s, Inches(0.9), Inches(0.85), Inches(4.5), Inches(0.5),
                    radius=0.5, line_color=C_CYAN, line_w=Pt(1))
    # 胶囊左侧小圆点
    add_oval(s, Inches(1.1), Inches(1.0), Inches(0.2), Inches(0.2), color=C_ORANGE)
    add_text(s, Inches(1.45), Inches(0.92), Inches(3.8), Inches(0.36),
             "昌吉州 AI 智能体创新应用大赛", size=13, color=C_CYAN, bold=True,
             align=PP_ALIGN.LEFT)
    # 主标题（超大）
    add_text(s, Inches(0.9), Inches(2.2), Inches(11), Inches(1.1),
             "AI 数字人工匠导师", size=58, color=C_TEXT_LT, bold=True)
    add_text(s, Inches(0.9), Inches(3.25), Inches(11), Inches(0.9),
             "智能体", size=58, color=C_TEXT_LT, bold=True)
    # 副标题（带色）
    add_text(s, Inches(0.9), Inches(4.3), Inches(11), Inches(0.55),
             "职业教育全栈元宇宙解决方案", size=26, color=C_CYAN, bold=True)
    # 一句话价值主张（带左侧色条）
    add_rect(s, Inches(0.9), Inches(5.15), Inches(0.06), Inches(0.4), color=C_ORANGE)
    add_text(s, Inches(1.1), Inches(5.15), Inches(10), Inches(0.4),
             "AI 大模型  ×  空间计算  ×  WebGL2.0  ——  打造沉浸式职教实训新范式",
             size=15, color=C_TEXT_DIM)
    # 底部三信息块
    infos = [
        ("定位", "依托昌吉州“东数西算”枢纽 · 国际融合算力中心", C_CYAN),
        ("方向", "AI + 职业教育数字化实训", C_ORANGE),
        ("形式", "路演汇报  /  2026", C_VIOLET),
    ]
    for i, (k, v, col) in enumerate(infos):
        y = Inches(6.0 + i * 0.42)
        add_text(s, Inches(0.9), y, Inches(0.8), Inches(0.35),
                 k, size=11, color=col, bold=True)
        add_text(s, Inches(1.7), y, Inches(9), Inches(0.35),
                 v, size=12, color=C_TEXT_DIM)


# ============ 2. 目录 ============
def slide_toc():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    # 左侧深色面板
    left = add_rect(s, 0, 0, Inches(4.5), SLIDE_H)
    set_gradient3(left, (0x05, 0x0A, 0x1A), (0x0A, 0x1F, 0x3D), (0x12, 0x2E, 0x5C), angle=90)
    deco_grid_bg(s, color=RGBColor(0x1A, 0x2E, 0x4D), alpha=10)
    # 左侧装饰
    deco_corner_rings(s, Inches(3.2), Inches(5.2), Inches(2.2), C_CYAN, line_w=Pt(1), alpha=30)
    deco_glow_dot(s, Inches(0.8), Inches(0.8), Inches(0.14), C_ORANGE, alpha=60)
    add_text(s, Inches(0.6), Inches(2.6), Inches(3.6), Inches(0.4),
             "CONTENTS", size=13, color=C_CYAN, bold=True, font=FONT_T)
    add_text(s, Inches(0.6), Inches(3.0), Inches(3.6), Inches(0.9),
             "目  录", size=46, color=C_TEXT_LT, bold=True)
    add_rect(s, Inches(0.6), Inches(4.0), Inches(1.5), Pt(3), color=C_ORANGE)
    add_text(s, Inches(0.6), Inches(4.25), Inches(3.6), Inches(0.4),
             "AI 数字人工匠导师智能体", size=13, color=C_TEXT_DIM)
    add_text(s, Inches(0.6), Inches(4.6), Inches(3.6), Inches(0.4),
             "路演汇报  ·  2026", size=13, color=C_TEXT_DIM)
    # 右侧目录
    items = [
        ("01", "行业痛点与立项背景", C_ORANGE),
        ("02", "项目目标与产品定位", C_CYAN),
        ("03", "总体技术架构", C_ROYAL),
        ("04", "三大核心技术引擎", C_VIOLET),
        ("05", "AI 智能体闭环能力", C_TEAL),
        ("06", "产品功能与应用场景", C_ORANGE),
        ("07", "成果与质效提升", C_CYAN),
        ("08", "经济与社会效益", C_ROYAL),
        ("09", "创新优势与推广路径", C_VIOLET),
        ("10", "发展规划与愿景", C_TEAL),
    ]
    col1 = items[:5]; col2 = items[5:]
    for i, (num, name, col) in enumerate(col1):
        y = Inches(0.85 + i * 1.18)
        # 编号
        add_text(s, Inches(5.0), y, Inches(0.85), Inches(0.5),
                 num, size=24, color=col, bold=True)
        # 名称
        add_text(s, Inches(5.85), y + Inches(0.1), Inches(3.0), Inches(0.5),
                 name, size=15, color=C_TEXT, bold=True)
        # 分隔线
        add_rect(s, Inches(5.0), y + Inches(0.6), Inches(3.5), Pt(0.5), color=C_LINE)
    for i, (num, name, col) in enumerate(col2):
        y = Inches(0.85 + i * 1.18)
        add_text(s, Inches(9.5), y, Inches(0.85), Inches(0.5),
                 num, size=24, color=col, bold=True)
        add_text(s, Inches(10.35), y + Inches(0.1), Inches(3.0), Inches(0.5),
                 name, size=15, color=C_TEXT, bold=True)
        add_rect(s, Inches(9.5), y + Inches(0.6), Inches(3.5), Pt(0.5), color=C_LINE)
    page_footer(s)


# ============ 章节分隔页 ============
def slide_section(num, title, subtitle, page_num):
    s = prs.slides.add_slide(BLANK)
    bg = add_rect(s, 0, 0, SLIDE_W, SLIDE_H)
    set_gradient3(bg, (0x05, 0x0A, 0x1A), (0x0A, 0x1F, 0x3D), (0x12, 0x2E, 0x5C), angle=120)
    deco_grid_bg(s, color=RGBColor(0x1A, 0x2E, 0x4D), alpha=10)
    # 右侧装饰圆环
    deco_corner_rings(s, Inches(9.8), Inches(1.5), Inches(3), C_CYAN, line_w=Pt(1.2), alpha=30)
    deco_corner_rings(s, Inches(10.3), Inches(2.0), Inches(2), C_ORANGE, line_w=Pt(0.8), alpha=35)
    # 大编号（带轮廓效果）
    add_text(s, Inches(0.8), Inches(1.8), Inches(5), Inches(2.5),
             num, size=160, color=RGBColor(0x14, 0x2E, 0x5C), bold=True)
    # 编号上方小标签
    add_text(s, Inches(0.85), Inches(1.5), Inches(3), Inches(0.3),
             "PART " + num, size=12, color=C_ORANGE, bold=True)
    # 竖线分隔
    add_rect(s, Inches(4.5), Inches(2.7), Pt(3), Inches(2), color=C_ORANGE)
    # 标题
    add_text(s, Inches(4.9), Inches(2.9), Inches(8), Inches(0.8),
             title, size=42, color=C_TEXT_LT, bold=True)
    # 副标题
    add_text(s, Inches(4.9), Inches(3.85), Inches(8), Inches(0.5),
             subtitle, size=16, color=C_CYAN)
    # 底部装饰线 + 页码
    add_rect(s, Inches(0.8), Inches(6.8), Inches(11.7), Pt(0.5), color=C_LINE_DARK)
    add_text(s, Inches(0.8), Inches(6.9), Inches(6), Inches(0.3),
             "AI 数字人工匠导师智能体  ·  路演", size=10, color=C_TEXT_DIM)
    add_text(s, Inches(11.8), Inches(6.9), Inches(1), Inches(0.3),
             f"{page_num:02d} / 19", size=10, color=C_TEXT_DIM, align=PP_ALIGN.RIGHT)


# ============ 3. 行业痛点 ============
def slide_pain_points():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 01  ·  行业痛点", "职业教育数字化实训的三大痛点", num=3)
    pains = [
        ("3D 实训内容创作门槛高",
         ["非专业教师难以独立完成 3D 课件制作",
          "高度依赖外部外包定制，周期长、成本高",
          "内容更新僵化，难以响应教学需求变化"],
         C_ORANGE, "01"),
        ("虚拟仿真局限于单机体验",
         ["缺乏高并发低延迟的多人协同机制",
          "交互刻板，无个性化指导与反馈",
          "难以开展跨地域远程协作实训"],
         C_CYAN, "02"),
        ("时空与硬件壁垒突出",
         ["新疆地域广阔，师资与资源分布不均",
          "偏远地区难以享受优质实训资源",
          "高端 VR/工作站硬件依赖，采购成本高"],
         C_VIOLET, "03"),
    ]
    for i, (title, lines, col, num) in enumerate(pains):
        x = Inches(0.85 + i * 4.1)
        y = Inches(1.95)
        w = Inches(3.85); h = Inches(4.5)
        # 卡片（带阴影）
        add_rrect(s, x, y, w, h, color=C_CARD, line_color=C_LINE, line_w=Pt(0.5),
                  radius=0.04, shadow=True)
        # 顶部渐变色块
        top = add_rrect(s, x, y, w, Inches(1.0), radius=0.04,
                        gradient=(col, (col[0]//3, col[1]//3, col[2]//3)), grad_angle=0)
        # 顶部色块只圆角顶部 - 用矩形覆盖底部
        add_rect(s, x, y + Inches(0.5), w, Inches(0.5), color=col)
        # 大编号
        add_text(s, x + Inches(0.3), y + Inches(0.15), Inches(1.2), Inches(0.7),
                 num, size=32, color=C_TEXT_LT, bold=True)
        # 标题
        add_text(s, x + Inches(1.3), y + Inches(0.28), w - Inches(1.5), Inches(0.5),
                 title, size=15, color=C_TEXT_LT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        # 内容
        add_bullets(s, x + Inches(0.35), y + Inches(1.3), w - Inches(0.7), h - Inches(1.5),
                    lines, size=13, bullet_color=col, line_spacing=1.55, gap=10)
    # 底部结语胶囊
    add_rrect(s, Inches(0.85), Inches(6.65), Inches(11.65), Inches(0.42),
              radius=0.5, color=RGBColor(0xE8, 0xF4, 0xFD))
    add_oval(s, Inches(1.05), Inches(6.78), Inches(0.16), Inches(0.16), color=C_ORANGE)
    add_text(s, Inches(1.35), Inches(6.68), Inches(11), Inches(0.36),
             "立项契机：将 AI 技术与空间计算、WebGL2.0 渲染融合，打造 AI 数字人工匠导师智能体",
             size=12, color=C_ROYAL, bold=True)
    page_footer(s)


# ============ 4. 项目目标 ============
def slide_goals():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 02  ·  项目目标", "四大项目目标", num=4)
    goals = [
        ("打造 AI 工匠导师智能体",
         "构建具备自主感知、推理与执行能力的 AI 数字人工匠导师，打造职业教育全栈元宇宙解决方案，破解职教数字化实训行业痛点。",
         C_CYAN, "01"),
        ("依托算力底座赋能实训",
         "依托昌吉州算力底座，实现实训教学的多人协同、零门槛 3D 课件创作、个性化伴学辅导，提升实训智能化、沉浸式与高效性。",
         C_ORANGE, "02"),
        ("消纳算力 · 区域标杆",
         "有效消纳昌吉州融合算力资源，推动算力与职教场景深度融合，打造区域数字职教标杆。",
         C_ROYAL, "03"),
        ("可复制模式 · 技能出海",
         "形成可复制、可推广的职教数字化模式，赋能昌吉州及全疆职业院校与企业实训，为“一带一路”技能出海提供数字化载体。",
         C_VIOLET, "04"),
    ]
    for i, (title, desc, col, num) in enumerate(goals):
        row = i // 2; ccol = i % 2
        x = Inches(0.85 + ccol * 6.1)
        y = Inches(1.95 + row * 2.55)
        w = Inches(5.85); h = Inches(2.3)
        # 卡片
        add_rrect(s, x, y, w, h, color=C_CARD, line_color=C_LINE, line_w=Pt(0.5),
                  radius=0.04, shadow=True)
        # 左侧色条
        add_rect(s, x, y, Inches(0.1), h, color=col)
        # 编号圆（带渐变）
        add_oval(s, x + Inches(0.4), y + Inches(0.4), Inches(0.75), Inches(0.75),
                 gradient=(col, (col[0]//2, col[1]//2, col[2]//2)), grad_angle=45)
        add_text(s, x + Inches(0.4), y + Inches(0.45), Inches(0.75), Inches(0.65),
                 num, size=20, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
        # 标题
        add_text(s, x + Inches(1.35), y + Inches(0.4), w - Inches(1.6), Inches(0.5),
                 title, size=17, color=C_TEXT, bold=True)
        # 装饰短线
        add_rect(s, x + Inches(1.35), y + Inches(0.95), Inches(0.6), Pt(2), color=col)
        # 描述
        add_text(s, x + Inches(1.35), y + Inches(1.15), w - Inches(1.6), h - Inches(1.3),
                 desc, size=12, color=C_TEXT_SUB, line_spacing=1.55)
    page_footer(s)


# ============ 5. 产品定位 ============
def slide_positioning():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 02  ·  产品定位", "一中心 · 三角色 · 全栈闭环", num=5)
    # 中心大圆（带渐变 + 光晕）
    cx = Inches(6.67); cy = Inches(4.3)
    center = add_oval(s, cx - Inches(1.5), cy - Inches(1.5), Inches(3), Inches(3),
                      gradient=(C_ROYAL, C_CYAN), grad_angle=45, glow=(C_CYAN, Pt(20), 45))
    # 中心环装饰
    deco_corner_rings(s, cx - Inches(1.7), cy - Inches(1.7), Inches(3.4),
                      C_CYAN, line_w=Pt(0.75), alpha=40)
    add_text(s, cx - Inches(1.5), cy - Inches(0.75), Inches(3), Inches(0.5),
             "AI 工匠导师", size=20, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, cx - Inches(1.5), cy - Inches(0.15), Inches(3), Inches(0.4),
             "智能体", size=20, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
    add_rect(s, cx - Inches(0.5), cy + Inches(0.3), Inches(1), Pt(1.5), color=C_ORANGE)
    add_text(s, cx - Inches(1.3), cy + Inches(0.5), Inches(2.6), Inches(0.3),
             "感知 · 推理 · 执行", size=11, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)
    # 三个角色卡片
    roles = [
        ("智能讲师", "AI 数字人授课", "3D 视角联动 · TTS 配音", Inches(1.1), Inches(1.95), C_ORANGE),
        ("智能助教", "个性化伴学辅导", "多模态问答 · 实时反馈", Inches(9.5), Inches(1.95), C_CYAN),
        ("智能教研员", "零门槛 3D 课件创作", "文本 → 大纲 → 3D 动态课件", Inches(1.1), Inches(5.0), C_VIOLET),
    ]
    for name, sub, desc, x, y, col in roles:
        w = Inches(3.0); h = Inches(1.65)
        add_rrect(s, x, y, w, h, color=C_CARD, line_color=col, line_w=Pt(1.5),
                  radius=0.06, shadow=True)
        # 顶部小色块
        add_rect(s, x + Inches(0.3), y + Inches(0.2), Inches(0.4), Pt(3), color=col)
        add_text(s, x + Inches(0.3), y + Inches(0.32), w - Inches(0.6), Inches(0.4),
                 name, size=16, color=col, bold=True)
        add_text(s, x + Inches(0.3), y + Inches(0.78), w - Inches(0.6), Inches(0.35),
                 sub, size=12, color=C_TEXT)
        add_text(s, x + Inches(0.3), y + Inches(1.15), w - Inches(0.6), Inches(0.35),
                 desc, size=10, color=C_TEXT_SUB)
    # 右下：全栈闭环
    add_rrect(s, Inches(9.5), Inches(5.0), Inches(3.0), Inches(1.65),
              gradient=(C_ROYAL, C_NAVY), grad_angle=45, radius=0.06, shadow=True)
    add_text(s, Inches(9.5), Inches(5.25), Inches(3.0), Inches(0.4),
             "全栈闭环", size=16, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(10.7), Inches(5.7), Inches(0.6), Pt(1.5), color=C_CYAN)
    add_text(s, Inches(9.7), Inches(5.85), Inches(2.6), Inches(0.7),
             "Web · Unity · VR\n三端一致播放",
             size=11, color=C_TEXT_DIM, align=PP_ALIGN.CENTER, line_spacing=1.4)
    # 连线（中心到四角色）- 虚线感用细线
    for (rx, ry, _, _) in [(Inches(2.6), Inches(2.3), None, None),
                            (Inches(10.0), Inches(2.3), None, None),
                            (Inches(2.6), Inches(5.4), None, None),
                            (Inches(10.0), Inches(5.4), None, None)]:
        pass  # 连线复杂，用装饰点代替
    page_footer(s)


# ============ 6. 总体技术架构 ============
def slide_architecture():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 03  ·  技术架构", "云-边-端 协同分布式架构", num=6)
    layers = [
        ("云  层", "昌吉智算中枢",
         ["AI 大模型推理\nDeepSeek / QwenVL", "3D 资产存储转码\nGLB 标准化", "微服务 + K8s\n弹性扩容", "数据湖治理\n全生命周期"],
         C_ROYAL, "CLOUD"),
        ("边缘层", "CDN + 边缘计算",
         ["3D 资产就近缓存\n边缘节点", "解决西部网络波动\n弱网降级", "预加载策略\n体验保障", "跨地域协同\n低延迟保障"],
         C_CYAN, "EDGE"),
        ("端  层", "WebGL2.0 + WebXR",
         ["轻量化渲染引擎\n浏览器原生", "全终端无缝接入\n无硬件依赖", "VR/PICO 沉浸\nWebXR 标准", "摆脱高配硬件\n普惠化"],
         C_ORANGE, "CLIENT"),
    ]
    for i, (name, sub, items, col, en) in enumerate(layers):
        y = Inches(1.9 + i * 1.7)
        # 左侧标签卡（带渐变 + 阴影）
        add_rrect(s, Inches(0.85), y, Inches(2.5), Inches(1.5),
                  gradient=(col, (col[0]//2, col[1]//2, col[2]//2)), grad_angle=45,
                  radius=0.06, shadow=True)
        add_text(s, Inches(0.85), y + Inches(0.2), Inches(2.5), Inches(0.45),
                 name, size=20, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(0.85), y + Inches(0.7), Inches(2.5), Inches(0.35),
                 sub, size=11, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)
        add_text(s, Inches(0.85), y + Inches(1.05), Inches(2.5), Inches(0.3),
                 en, size=9, color=RGBColor(0x80, 0xA0, 0xC0), align=PP_ALIGN.CENTER, font=FONT_T)
        # 右侧内容卡
        add_rrect(s, Inches(3.55), y, Inches(8.95), Inches(1.5),
                  color=C_CARD, line_color=C_LINE, line_w=Pt(0.5), radius=0.04, shadow=True)
        # 四列要点
        for j, item in enumerate(items):
            ix = Inches(3.8 + j * 2.18)
            # 小圆点
            add_oval(s, ix, y + Inches(0.25), Inches(0.12), Inches(0.12), color=col)
            add_text(s, ix + Inches(0.22), y + Inches(0.18), Inches(1.9), Inches(1.2),
                     item, size=11, color=C_TEXT, line_spacing=1.4)
        # 层间连接小箭头
        if i < 2:
            add_text(s, Inches(1.9), y + Inches(1.5), Inches(0.4), Inches(0.2),
                     "▼", size=10, color=C_TEXT_SUB, align=PP_ALIGN.CENTER)
    # 底部说明
    add_rrect(s, Inches(0.85), Inches(7.0), Inches(11.65), Inches(0.32),
              radius=0.5, color=RGBColor(0xE8, 0xF4, 0xFD))
    add_text(s, Inches(0.85), Inches(7.02), Inches(11.65), Inches(0.28),
             "高可用  ·  高并发  ·  高沉浸  ——  AI 职教元宇宙技术生态",
             size=11, color=C_ROYAL, bold=True, align=PP_ALIGN.CENTER)
    page_footer(s)


# ============ 7. 核心引擎 1：AI 智能 3D 编辑器 ============
def slide_engine1():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 04  ·  核心引擎 01", "AI 智能 3D 编辑器", num=7)
    # 左侧：能力列表
    add_rrect(s, Inches(0.85), Inches(1.95), Inches(6.1), Inches(4.65),
              color=C_CARD, line_color=C_LINE, line_w=Pt(0.5), radius=0.03, shadow=True)
    add_text(s, Inches(1.15), Inches(2.15), Inches(5.5), Inches(0.4),
             "核心能力", size=16, color=C_TEXT, bold=True)
    add_rect(s, Inches(1.15), Inches(2.6), Inches(0.7), Pt(2), color=C_ORANGE)
    caps = [
        ("多格式 3D 模型解析", "FBX / OBJ / STL / GLB / STEP，自动转码 GLB"),
        ("AI 自动生成课件", "“文本描述 → 大纲 → 3D 动态课件” 一键生成"),
        ("AI 智能整理", "QwenVL 视觉大模型自动识别重命名部件、分组结构树"),
        ("关键帧动画编辑", "相机/显隐/变换三轨道，精确控制每帧状态"),
        ("一键演示配置", "相机关键帧 · 部件高亮 · 语音解说 · 标注联动"),
    ]
    for i, (t, d) in enumerate(caps):
        y = Inches(2.85 + i * 0.74)
        # 序号圆
        add_oval(s, Inches(1.15), y + Inches(0.05), Inches(0.32), Inches(0.32),
                 gradient=(C_ORANGE, C_AMBER), grad_angle=45)
        add_text(s, Inches(1.15), y + Inches(0.07), Inches(0.32), Inches(0.28),
                 f"{i+1}", size=11, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
        # 标题
        add_text(s, Inches(1.6), y + Inches(0.05), Inches(2.3), Inches(0.3),
                 t, size=13, color=C_TEXT, bold=True)
        # 描述
        add_text(s, Inches(3.95), y + Inches(0.05), Inches(2.9), Inches(0.6),
                 d, size=10, color=C_TEXT_SUB, line_spacing=1.3)
    # 右侧：工作流（深色卡片）
    add_rrect(s, Inches(7.15), Inches(1.95), Inches(5.35), Inches(4.65),
              gradient=(C_INK, C_NAVY), grad_angle=45, radius=0.03, shadow=True)
    # 右上装饰
    deco_corner_rings(s, Inches(11.8), Inches(2.15), Inches(0.8), C_CYAN, line_w=Pt(0.75), alpha=50)
    add_text(s, Inches(7.45), Inches(2.15), Inches(5), Inches(0.4),
             "AI 课件生成工作流", size=16, color=C_CYAN, bold=True)
    add_rect(s, Inches(7.45), Inches(2.6), Inches(0.7), Pt(2), color=C_ORANGE)
    flow = [
        "教师输入文本描述 / 教学目标",
        "AI 大模型生成课程大纲",
        "自动配置相机关键帧 · 部件高亮",
        "生成 3D 动态课件 + 讲稿",
        "TTS 批量配音 → 发布",
    ]
    for i, step in enumerate(flow):
        y = Inches(2.95 + i * 0.72)
        # 步骤胶囊
        add_rrect(s, Inches(7.55), y, Inches(4.6), Inches(0.5),
                  color=C_CARD_MID, line_color=C_CYAN, line_w=Pt(0.5), radius=0.2)
        # 序号
        add_text(s, Inches(7.7), y + Inches(0.08), Inches(0.4), Inches(0.32),
                 f"{i+1}", size=12, color=C_ORANGE, bold=True)
        # 步骤文本
        add_text(s, Inches(8.1), y + Inches(0.1), Inches(4), Inches(0.3),
                 step, size=12, color=C_TEXT_LT)
        # 箭头
        if i < len(flow) - 1:
            add_text(s, Inches(9.5), y + Inches(0.5), Inches(0.6), Inches(0.25),
                     "▼", size=12, color=C_CYAN, align=PP_ALIGN.CENTER)
    page_footer(s)


# ============ 8. 核心引擎 2：超低延迟通信 ============
def slide_engine2():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 04  ·  核心引擎 02", "超低延迟通信协同", num=8)
    # 四个数据卡（深色玻璃风格）
    stats = [
        ("<100", "ms", "音频通信延迟", "WebRTC", C_ORANGE),
        ("<50", "ms", "白板同步延迟", "WebSocket", C_CYAN),
        ("50", "人", "同频在线协同", "高并发", C_VIOLET),
        ("6", "DoF", "动作捕捉重定向", "沉浸式", C_TEAL),
    ]
    for i, (num, unit, label, tag, col) in enumerate(stats):
        x = Inches(0.85 + i * 3.05)
        y = Inches(1.95)
        w = Inches(2.85); h = Inches(1.55)
        # 深色卡片 + 渐变
        add_rrect(s, x, y, w, h, gradient=(C_INK, C_NAVY), grad_angle=45,
                  radius=0.05, shadow=True)
        # 顶部色条
        add_rect(s, x + Inches(0.3), y + Inches(0.15), Inches(0.5), Pt(3), color=col)
        # 大数字
        add_text(s, x + Inches(0.25), y + Inches(0.35), w - Inches(0.5), Inches(0.6),
                 num, size=36, color=col, bold=True, align=PP_ALIGN.LEFT)
        # 单位
        add_text(s, x + Inches(1.7), y + Inches(0.55), Inches(0.8), Inches(0.4),
                 unit, size=16, color=C_TEXT_DIM, align=PP_ALIGN.LEFT)
        # 标签
        add_text(s, x + Inches(0.3), y + Inches(1.05), w - Inches(0.6), Inches(0.3),
                 label, size=12, color=C_TEXT_LT, bold=True)
        # 技术标签
        add_text(s, x + Inches(0.3), y + Inches(1.3), w - Inches(0.6), Inches(0.2),
                 tag, size=9, color=col, font=FONT_T)
    # 下方三能力卡
    cards = [
        ("WebRTC 音频通信", C_ORANGE,
         ["高清语音多人通话", "3D 空间音频定位", "智能混音降噪", "延迟控制在 100ms 内"]),
        ("WebSocket 实时同步", C_CYAN,
         ["虚拟白板 50ms 同步", "教师授权学生操作", "同步 3D 模型操作", "多人实时协作标记"]),
        ("沉浸式协同体验", C_VIOLET,
         ["3D 空间音频沉浸感", "视角强制同步", "VR 飞屏技术", "跨地域协同实训"]),
    ]
    for i, (title, col, lines) in enumerate(cards):
        x = Inches(0.85 + i * 4.1)
        y = Inches(3.8)
        w = Inches(3.85); h = Inches(2.9)
        add_rrect(s, x, y, w, h, color=C_CARD, line_color=C_LINE, line_w=Pt(0.5),
                  radius=0.04, shadow=True)
        # 标题区
        add_rect(s, x + Inches(0.3), y + Inches(0.25), Inches(0.5), Pt(3), color=col)
        add_text(s, x + Inches(0.3), y + Inches(0.4), w - Inches(0.6), Inches(0.4),
                 title, size=15, color=col, bold=True)
        # 内容
        add_bullets(s, x + Inches(0.3), y + Inches(0.95), w - Inches(0.6), h - Inches(1.1),
                    lines, size=12, bullet_color=col, line_spacing=1.6, gap=8)
    page_footer(s)


# ============ 9. 核心引擎 3：高并发渲染 ============
def slide_engine3():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 04  ·  核心引擎 03", "高并发渲染优化", num=9)
    # 左侧：技术手段
    add_text(s, Inches(0.85), Inches(1.9), Inches(6), Inches(0.4),
             "关键技术手段", size=16, color=C_TEXT, bold=True)
    techs = [
        ("动态 LOD", "根据距离动态切换模型精度，远距简化、近距精细", C_ORANGE),
        ("遮挡剔除", "智能剔除被遮挡几何体，减少不必要的绘制调用", C_CYAN),
        ("实例化渲染", "相同几何体单次绘制批量渲染，大幅降低 GPU 负担", C_ROYAL),
        ("6DoF 动作捕捉", "头盔/手柄跟踪 + 位置算法驱动全身姿态", C_VIOLET),
        ("骨骼重定向", "统一骨骼骨架映射，支持多 Avatar 动作同步", C_TEAL),
    ]
    for i, (t, d, col) in enumerate(techs):
        y = Inches(2.4 + i * 0.82)
        add_rrect(s, Inches(0.85), y, Inches(6.3), Inches(0.72),
                  color=C_CARD, line_color=C_LINE, line_w=Pt(0.5), radius=0.06, shadow=True)
        add_rect(s, Inches(0.85), y, Inches(0.1), Inches(0.72), color=col)
        # 技术名
        add_text(s, Inches(1.15), y + Inches(0.12), Inches(2.0), Inches(0.3),
                 t, size=13, color=col, bold=True)
        # 描述
        add_text(s, Inches(1.15), y + Inches(0.42), Inches(5.0), Inches(0.3),
                 d, size=10, color=C_TEXT_SUB)
    # 右侧：性能指标（深色面板）
    add_rrect(s, Inches(7.35), Inches(1.95), Inches(5.15), Inches(4.65),
              gradient=(C_INK, C_NAVY), grad_angle=45, radius=0.03, shadow=True)
    deco_corner_rings(s, Inches(11.8), Inches(2.1), Inches(0.8), C_CYAN, line_w=Pt(0.75), alpha=40)
    add_text(s, Inches(7.65), Inches(2.15), Inches(4.7), Inches(0.4),
             "性能指标达成", size=16, color=C_CYAN, bold=True)
    add_rect(s, Inches(7.65), Inches(2.6), Inches(0.7), Pt(2), color=C_ORANGE)
    metrics = [
        ("50", "人并发同频", C_ORANGE),
        ("100万+", "面片级流畅渲染", C_CYAN),
        ("60", "FPS Web 端", C_ROYAL),
        ("90", "FPS VR 端", C_VIOLET),
        ("≤200", "ms 操作响应", C_TEAL),
        ("2", "Mbps 弱网可用", C_ORANGE),
    ]
    for i, (num, label, col) in enumerate(metrics):
        row = i // 2; ccol = i % 2
        x = Inches(7.75 + ccol * 2.4)
        y = Inches(2.95 + row * 1.15)
        # 小卡
        add_rrect(s, x, y, Inches(2.2), Inches(1.0),
                  color=C_CARD_MID, line_color=C_LINE_DARK, line_w=Pt(0.5), radius=0.08)
        add_text(s, x, y + Inches(0.15), Inches(2.2), Inches(0.45),
                 num, size=24, color=col, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.62), Inches(2.2), Inches(0.3),
                 label, size=10, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)
    page_footer(s)


# ============ 10. AI 智能体闭环 ============
def slide_ai_agent():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 05  ·  AI 智能体", "“感知 - 推理 - 执行” 闭环智能体", num=10)
    stages = [
        ("感  知", "ASR · CV · 多模态",
         "语音 / 文字 / 图片 / 截图\n多模态输入理解", C_CYAN, Inches(1.3)),
        ("推  理", "DeepSeek 大模型",
         "课程大纲生成\n知识问答 · 个性化辅导", C_ORANGE, Inches(5.4)),
        ("执  行", "TTS · 3D 联动 · 数字人",
         "数字人讲解演示\n3D 视角切换 · 部件高亮", C_VIOLET, Inches(9.5)),
    ]
    for name, tech, desc, col, x in stages:
        # 大圆（渐变 + 光晕）
        add_oval(s, x, Inches(2.15), Inches(2.8), Inches(2.8),
                 gradient=(col, (col[0]//3, col[1]//3, col[2]//3)), grad_angle=45,
                 glow=(col, Pt(22), 45))
        # 外环
        deco_corner_rings(s, x - Inches(0.15), Inches(2.0), Inches(3.1),
                          col, line_w=Pt(0.75), alpha=35)
        add_text(s, x, Inches(2.95), Inches(2.8), Inches(0.5),
                 name, size=24, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(3.5), Inches(2.8), Inches(0.35),
                 tech, size=11, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)
        # 下方描述卡
        add_rrect(s, x - Inches(0.2), Inches(5.15), Inches(3.2), Inches(1.0),
                  color=C_CARD, line_color=C_LINE, line_w=Pt(0.5), radius=0.06, shadow=True)
        add_text(s, x - Inches(0.2), Inches(5.3), Inches(3.2), Inches(0.8),
                 desc, size=12, color=C_TEXT, align=PP_ALIGN.CENTER, line_spacing=1.5)
    # 箭头
    for ax in [Inches(4.1), Inches(8.2)]:
        add_text(s, ax, Inches(3.25), Inches(1.3), Inches(0.6),
                 "→", size=44, color=C_ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # 回环提示
    add_rrect(s, Inches(0.85), Inches(6.35), Inches(11.65), Inches(0.42),
              radius=0.5, gradient=(C_ROYAL, C_NAVY), grad_angle=0)
    add_text(s, Inches(0.85), Inches(6.38), Inches(11.65), Inches(0.36),
             "↻  闭环反馈：执行结果回流至感知层，持续优化教学策略",
             size=13, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
    # 合规胶囊
    add_rrect(s, Inches(0.85), Inches(6.88), Inches(11.65), Inches(0.32),
              radius=0.5, color=RGBColor(0xE8, 0xF4, 0xFD))
    add_text(s, Inches(0.85), Inches(6.9), Inches(11.65), Inches(0.28),
             "算法伦理围栏  ·  人机协同管控  ·  数据全生命周期治理  ——  保障系统安全合规",
             size=10, color=C_ROYAL, align=PP_ALIGN.CENTER)
    page_footer(s)


# ============ 11. 产品功能矩阵 ============
def slide_features():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 06  ·  产品功能", "六大核心产品模块", num=11)
    feats = [
        ("元宇宙教学大厅", "Unity / PICO", "多人 VR 协同教学\nAvatar · 白板 · 飞屏", C_ORANGE, "01"),
        ("AI 智能 3D 编辑器", "Web · Three.js", "零门槛 3D 课件创作\nAI 整理 · 关键帧动画", C_CYAN, "02"),
        ("AI 课件制作", "DeepSeek 大模型", "一键生成授课大纲\n考题 · TTS 配音", C_ROYAL, "03"),
        ("公开播放器", "Web · WebXR", "学习/探索/答题三模式\nVR 沉浸式体验", C_VIOLET, "04"),
        ("资源与用户管理", "Next.js + Express", "多角色权限 · 配额\n批量导入 · 资源转码", C_TEAL, "05"),
        ("数据分析看板", "MongoDB 聚合", "成绩统计 · 学习进度\n趋势图 · 多级穿透", C_MAGENTA, "06"),
    ]
    for i, (name, tech, desc, col, num) in enumerate(feats):
        row = i // 3; ccol = i % 3
        x = Inches(0.85 + ccol * 4.1)
        y = Inches(1.95 + row * 2.55)
        w = Inches(3.85); h = Inches(2.35)
        # 卡片
        add_rrect(s, x, y, w, h, color=C_CARD, line_color=C_LINE, line_w=Pt(0.5),
                  radius=0.04, shadow=True)
        # 顶部渐变条
        top = add_rrect(s, x, y, w, Inches(0.55), radius=0.04,
                        gradient=(col, (col[0]//2, col[1]//2, col[2]//2)), grad_angle=0)
        # 覆盖底部圆角
        add_rect(s, x, y + Inches(0.3), w, Inches(0.25), color=col)
        # 编号
        add_text(s, x + Inches(0.3), y + Inches(0.1), Inches(0.6), Inches(0.35),
                 num, size=14, color=C_TEXT_LT, bold=True)
        # 名称
        add_text(s, x + Inches(0.85), y + Inches(0.12), w - Inches(1.0), Inches(0.35),
                 name, size=14, color=C_TEXT_LT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        # 技术标签
        add_rrect(s, x + Inches(0.3), y + Inches(0.75), Inches(2.2), Inches(0.3),
                  radius=0.5, color=RGBColor(0xF0, 0xF4, 0xFA))
        add_text(s, x + Inches(0.3), y + Inches(0.78), Inches(2.2), Inches(0.24),
                 tech, size=10, color=col, bold=True, align=PP_ALIGN.CENTER, font=FONT_T)
        # 描述
        add_text(s, x + Inches(0.3), y + Inches(1.25), w - Inches(0.6), Inches(1.0),
                 desc, size=12, color=C_TEXT, line_spacing=1.55)
    page_footer(s)


# ============ 12. 应用场景 ============
def slide_scenarios():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 06  ·  应用场景", "多行业 · 多专业 · 多形态", num=12)
    scenes = [
        ("高危特种作业实训", "云端 1:1 还原工业现场，零风险虚拟实操，体验违规操作后果", C_ORANGE, "01"),
        ("装配式建筑专业", "节点构造 · 结构展示 · 装配工艺流程三维可视化教学", C_CYAN, "02"),
        ("自动化产线专业", "工业机器人 · 产线调试 · 设备维护三维仿真实训", C_ROYAL, "03"),
        ("水利水电专业", "枢纽工程 · 设备原理 · 施工工艺沉浸式教学", C_VIOLET, "04"),
        ("企业技能内训", "降低企业内训成本，标准化岗位技能培训与考核", C_TEAL, "05"),
        ("“一带一路”出海", "多语种实训平台，输出中国职业技能标准", C_MAGENTA, "06"),
    ]
    for i, (name, desc, col, num) in enumerate(scenes):
        row = i // 3; ccol = i % 3
        x = Inches(0.85 + ccol * 4.1)
        y = Inches(1.95 + row * 2.55)
        w = Inches(3.85); h = Inches(2.35)
        add_rrect(s, x, y, w, h, color=C_CARD, line_color=C_LINE, line_w=Pt(0.5),
                  radius=0.04, shadow=True)
        # 编号圆
        add_oval(s, x + Inches(0.35), y + Inches(0.4), Inches(0.85), Inches(0.85),
                 gradient=(col, (col[0]//2, col[1]//2, col[2]//2)), grad_angle=45)
        add_text(s, x + Inches(0.35), y + Inches(0.48), Inches(0.85), Inches(0.65),
                 num, size=22, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
        # 标题
        add_text(s, x + Inches(1.35), y + Inches(0.5), w - Inches(1.6), Inches(0.6),
                 name, size=15, color=C_TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        # 分隔线
        add_rect(s, x + Inches(0.35), y + Inches(1.4), w - Inches(0.7), Pt(0.5), color=C_LINE)
        # 描述
        add_text(s, x + Inches(0.35), y + Inches(1.55), w - Inches(0.7), Inches(0.75),
                 desc, size=11, color=C_TEXT_SUB, line_spacing=1.55)
    page_footer(s)


# ============ 13. 痛点解决成果 ============
def slide_pain_solution():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 07  ·  成果", "痛点解决成果", num=13)
    items = [
        ("破解高危特种作业实训难题",
         "云端 1:1 还原工业现场，实现零风险虚拟实操，让学员在安全环境中体验违规操作后果，解决实体实训高消耗、高风险、难复现问题。",
         C_ORANGE, "01"),
        ("打破 3D 课件创作壁垒",
         "AI 智能编辑器让非专业教师通过文本描述快速生成/修改 3D 课件，终结外部外包定制的僵化模式，释放职教教研产能。",
         C_CYAN, "02"),
        ("消解时空与硬件壁垒",
         "全终端无缝接入模式，解决新疆地域广阔、师资资源不均问题，让偏远地区师生也能享受优质实训资源，摆脱高端硬件依赖。",
         C_VIOLET, "03"),
    ]
    for i, (title, desc, col, num) in enumerate(items):
        y = Inches(2.0 + i * 1.6)
        # 卡片
        add_rrect(s, Inches(0.85), y, Inches(11.65), Inches(1.4),
                  color=C_CARD, line_color=C_LINE, line_w=Pt(0.5), radius=0.04, shadow=True)
        # 左侧色条
        add_rect(s, Inches(0.85), y, Inches(0.15), Inches(1.4), color=col)
        # 大编号
        add_oval(s, Inches(1.2), y + Inches(0.3), Inches(0.8), Inches(0.8),
                 gradient=(col, (col[0]//2, col[1]//2, col[2]//2)), grad_angle=45)
        add_text(s, Inches(1.2), y + Inches(0.38), Inches(0.8), Inches(0.65),
                 num, size=22, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
        # 标题
        add_text(s, Inches(2.3), y + Inches(0.28), Inches(9.5), Inches(0.4),
                 title, size=17, color=C_TEXT, bold=True)
        add_rect(s, Inches(2.3), y + Inches(0.72), Inches(0.6), Pt(2), color=col)
        # 描述
        add_text(s, Inches(2.3), y + Inches(0.85), Inches(9.8), Inches(0.5),
                 desc, size=12, color=C_TEXT_SUB, line_spacing=1.5)
    page_footer(s)


# ============ 14. 质效提升成果 ============
def slide_quality():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 07  ·  成果", "质效提升成果  ·  降本 · 提效 · 增质", num=14)
    cards = [
        ("降  本", C_ORANGE, [
            ("数周 → 分钟级", "3D 课件制作周期"),
            ("零边际成本", "轻量化生产模式"),
            ("大幅降低", "时间与资金成本"),
        ]),
        ("提  效", C_CYAN, [
            ("50 人", "单实训大厅高并发"),
            ("量级跃升", "核心资源利用率"),
            ("跨地域协同", "与线下无异的交互"),
        ]),
        ("增  质", C_VIOLET, [
            ("沉浸式体验", "肌肉记忆 · 空间感"),
            ("技能转化率", "虚拟 → 真实岗位"),
            ("多维信息融合", "视角同步 · VR 飞屏"),
        ]),
    ]
    for i, (title, col, metrics) in enumerate(cards):
        x = Inches(0.85 + i * 4.1)
        y = Inches(1.95)
        w = Inches(3.85); h = Inches(4.85)
        # 卡片
        add_rrect(s, x, y, w, h, color=C_CARD, line_color=C_LINE, line_w=Pt(0.5),
                  radius=0.03, shadow=True)
        # 顶部渐变标题
        top = add_rrect(s, x, y, w, Inches(0.85), radius=0.03,
                        gradient=(col, (col[0]//2, col[1]//2, col[2]//2)), grad_angle=0)
        add_rect(s, x, y + Inches(0.5), w, Inches(0.35), color=col)
        add_text(s, x, y + Inches(0.18), w, Inches(0.5),
                 title, size=24, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
        # 三个指标
        for j, (val, label) in enumerate(metrics):
            my = y + Inches(1.15 + j * 1.25)
            # 左侧色条
            add_rect(s, x + Inches(0.35), my, Inches(0.08), Inches(0.85), color=col)
            # 数值
            add_text(s, x + Inches(0.55), my, w - Inches(0.9), Inches(0.5),
                     val, size=18, color=col, bold=True)
            # 标签
            add_text(s, x + Inches(0.55), my + Inches(0.55), w - Inches(0.9), Inches(0.3),
                     label, size=11, color=C_TEXT_SUB)
            # 分隔线
            if j < len(metrics) - 1:
                add_rect(s, x + Inches(0.35), my + Inches(1.05), w - Inches(0.7), Pt(0.5), color=C_LINE)
    page_footer(s)


# ============ 15. 经济与社会效益 ============
def slide_benefit():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 08  ·  效益", "经济与社会效益", num=15)
    items = [
        ("经济效益", C_ORANGE, [
            "高效消纳昌吉州算力资源，推动“东数西算”价值变现",
            "降低本地职教财政硬件采购成本与企业内训成本",
            "吸引 AI 职教上下游产业集聚，打造“算力+AI职教”产业链",
        ]),
        ("社会效益", C_CYAN, [
            "填平地域与师资鸿沟，实现职教资源全域普惠",
            "提升新疆青年数字素养与工业技能",
            "依托“一带一路”核心区，打造多语种实训平台",
            "成为中国职业技能标准出海的数字化桥头堡",
        ]),
        ("产业与示范价值", C_VIOLET, [
            "形成“政校企”协同的数字职教生态",
            "打造国家级虚拟仿真实训示范基地",
            "模式可快速复制至全疆乃至西北各地",
            "成为 AI 赋能职业教育升级的标杆案例",
        ]),
    ]
    for i, (title, col, lines) in enumerate(items):
        x = Inches(0.85 + i * 4.1)
        y = Inches(1.95)
        w = Inches(3.85); h = Inches(4.85)
        add_rrect(s, x, y, w, h, color=C_CARD, line_color=C_LINE, line_w=Pt(0.5),
                  radius=0.03, shadow=True)
        # 顶部渐变
        top = add_rrect(s, x, y, w, Inches(0.85), radius=0.03,
                        gradient=(col, (col[0]//2, col[1]//2, col[2]//2)), grad_angle=0)
        add_rect(s, x, y + Inches(0.5), w, Inches(0.35), color=col)
        add_text(s, x, y + Inches(0.2), w, Inches(0.5),
                 title, size=19, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
        # 内容
        add_bullets(s, x + Inches(0.35), y + Inches(1.15), w - Inches(0.7), h - Inches(1.3),
                    lines, size=12, bullet_color=col, line_spacing=1.7, gap=10)
    page_footer(s)


# ============ 16. 创新优势 ============
def slide_innovation():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 09  ·  优势", "五大创新优势", num=16)
    advs = [
        ("AI 原生", "AI 大模型深度融入 3D 创作与教学全链路，非“3D + AI 外挂”式拼凑", C_ORANGE),
        ("全栈闭环", "Web · Unity · VR 三端数据格式统一，创作-发布-学习-考核闭环", C_CYAN),
        ("算力消纳", "精准契合昌吉州“东数西算”定位，算力与职教场景深度融合", C_ROYAL),
        ("低门槛", "非专业教师可零门槛创作 3D 课件，全终端浏览器接入无硬件依赖", C_VIOLET),
        ("可复制", "模式标准化，可快速复制至全疆、西北乃至“一带一路”沿线国家", C_TEAL),
    ]
    for i, (title, desc, col) in enumerate(advs):
        y = Inches(1.95 + i * 1.0)
        # 卡片
        add_rrect(s, Inches(0.85), y, Inches(11.65), Inches(0.85),
                  color=C_CARD, line_color=C_LINE, line_w=Pt(0.5), radius=0.06, shadow=True)
        add_rect(s, Inches(0.85), y, Inches(0.12), Inches(0.85), color=col)
        # 编号圆
        add_oval(s, Inches(1.2), y + Inches(0.18), Inches(0.5), Inches(0.5),
                 gradient=(col, (col[0]//2, col[1]//2, col[2]//2)), grad_angle=45)
        add_text(s, Inches(1.2), y + Inches(0.22), Inches(0.5), Inches(0.4),
                 f"{i+1}", size=15, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
        # 标题
        add_text(s, Inches(2.0), y + Inches(0.12), Inches(2.5), Inches(0.35),
                 title, size=16, color=col, bold=True)
        # 描述
        add_text(s, Inches(2.0), y + Inches(0.48), Inches(9.8), Inches(0.35),
                 desc, size=12, color=C_TEXT_SUB)
    page_footer(s)


# ============ 17. 推广路径 ============
def slide_promotion():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 09  ·  推广", "三步走推广路径", num=17)
    steps = [
        ("第一阶段", "标杆示范", "昌吉州本地试点",
         ["遴选昌吉州职业院校试点部署", "打造区域数字职教标杆案例", "验证商业模式与教学效果"],
         C_CYAN, "01"),
        ("第二阶段", "区域复制", "全疆及西北推广",
         ["复制至全疆职业院校与企业实训", "拓展至西北五省职教市场", "建立“政校企”协同生态"],
         C_ORANGE, "02"),
        ("第三阶段", "出海输出", "“一带一路”桥头堡",
         ["多语种实训平台建设", "输出中国职业技能标准", "技术软实力国际化"],
         C_VIOLET, "03"),
    ]
    for i, (phase, keyword, sub, lines, col, num) in enumerate(steps):
        x = Inches(0.85 + i * 4.1)
        y = Inches(1.95)
        w = Inches(3.85); h = Inches(4.85)
        # 卡片
        add_rrect(s, x, y, w, h, color=C_CARD, line_color=C_LINE, line_w=Pt(0.5),
                  radius=0.03, shadow=True)
        # 顶部渐变
        top = add_rrect(s, x, y, w, Inches(1.3), radius=0.03,
                        gradient=(col, (col[0]//2, col[1]//2, col[2]//2)), grad_angle=0)
        add_rect(s, x, y + Inches(0.95), w, Inches(0.35), color=col)
        # 阶段编号
        add_text(s, x + Inches(0.3), y + Inches(0.18), Inches(1), Inches(0.35),
                 num, size=32, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
        # 阶段
        add_text(s, x + Inches(1.3), y + Inches(0.22), w - Inches(1.5), Inches(0.3),
                 phase, size=11, color=C_TEXT_DIM, align=PP_ALIGN.RIGHT)
        # 关键词
        add_text(s, x, y + Inches(0.62), w, Inches(0.5),
                 keyword, size=24, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
        # 副标题
        add_text(s, x, y + Inches(1.05), w, Inches(0.25),
                 sub, size=11, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)
        # 内容
        add_bullets(s, x + Inches(0.35), y + Inches(1.65), w - Inches(0.7), h - Inches(1.85),
                    lines, size=13, bullet_color=col, line_spacing=1.8, gap=12)
        # 阶段间箭头
        if i < 2:
            add_text(s, x + w - Inches(0.15), y + Inches(2.2), Inches(0.4), Inches(0.6),
                     "▶", size=22, color=C_ORANGE, bold=True, align=PP_ALIGN.CENTER)
    page_footer(s)


# ============ 18. 发展规划 ============
def slide_roadmap():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, color=C_PAPER)
    page_header(s, "PART 10  ·  愿景", "发展规划路线图", num=18)
    # 时间轴主线（渐变）
    line = add_rect(s, Inches(1.2), Inches(4.0), Inches(11), Pt(4))
    set_gradient(line, (0x00, 0xD4, 0xFF), (0xEC, 0x48, 0x99), angle=0)
    phases = [
        ("2026 H2", "产品深化", "完善 AI 智能体能力\n拓展专业场景库", C_CYAN),
        ("2027", "规模推广", "全疆职校规模化落地\n企业内训市场拓展", C_ORANGE),
        ("2028", "生态构建", "“政校企”协同生态\n国家级示范基地", C_VIOLET),
        ("2029+", "出海输出", "多语种平台\n“一带一路”技能出海", C_MAGENTA),
    ]
    for i, (time, title, desc, col) in enumerate(phases):
        cx = Inches(1.7 + i * 2.95)
        # 节点大圆（带光晕）
        add_oval(s, cx - Inches(0.3), Inches(3.7), Inches(0.6), Inches(0.6),
                 gradient=(col, (col[0]//2, col[1]//2, col[2]//2)), grad_angle=45,
                 glow=(col, Pt(15), 50))
        # 时间标签（上方）
        add_rrect(s, cx - Inches(0.85), Inches(2.35), Inches(1.7), Inches(0.4),
                  radius=0.5, color=C_CARD, line_color=col, line_w=Pt(1))
        add_text(s, cx - Inches(0.85), Inches(2.4), Inches(1.7), Inches(0.3),
                 time, size=14, color=col, bold=True, align=PP_ALIGN.CENTER, font=FONT_T)
        # 标题（上方）
        add_text(s, cx - Inches(1.2), Inches(2.95), Inches(2.4), Inches(0.4),
                 title, size=16, color=C_TEXT, bold=True, align=PP_ALIGN.CENTER)
        # 卡片（下方）
        add_rrect(s, cx - Inches(1.25), Inches(4.5), Inches(2.5), Inches(1.85),
                  color=C_CARD, line_color=col, line_w=Pt(1), radius=0.05, shadow=True)
        # 顶部色条
        add_rect(s, cx - Inches(1.25), Inches(4.5), Inches(2.5), Pt(3), color=col)
        add_text(s, cx - Inches(1.15), Inches(4.75), Inches(2.3), Inches(1.5),
                 desc, size=12, color=C_TEXT_SUB, align=PP_ALIGN.CENTER, line_spacing=1.7)
    # 底部愿景条
    add_rrect(s, Inches(0.85), Inches(6.6), Inches(11.65), Inches(0.5),
              gradient=(C_INK, C_NAVY), grad_angle=0, radius=0.1, shadow=True)
    add_text(s, Inches(0.85), Inches(6.65), Inches(11.65), Inches(0.4),
             "愿景：成为 AI 赋能职业教育升级的标杆，让每一个学习者都有专属的 AI 工匠导师",
             size=13, color=C_CYAN, bold=True, align=PP_ALIGN.CENTER)
    page_footer(s)


# ============ 19. 结尾页 ============
def slide_ending():
    s = prs.slides.add_slide(BLANK)
    bg = add_rect(s, 0, 0, SLIDE_W, SLIDE_H)
    set_gradient3(bg, (0x05, 0x0A, 0x1A), (0x0A, 0x1F, 0x3D), (0x12, 0x2E, 0x5C), angle=120)
    deco_grid_bg(s, color=RGBColor(0x1A, 0x2E, 0x4D), alpha=10)
    # 装饰圆环组
    deco_corner_rings(s, Inches(9.5), Inches(0.8), Inches(3.2), C_CYAN, line_w=Pt(1.2), alpha=35)
    deco_corner_rings(s, Inches(10.0), Inches(1.3), Inches(2.2), C_ORANGE, line_w=Pt(0.8), alpha=40)
    deco_corner_rings(s, Inches(10.4), Inches(1.7), Inches(1.4), C_VIOLET, line_w=Pt(0.6), alpha=35)
    # 光晕点
    deco_glow_dot(s, Inches(2.0), Inches(1.5), Inches(0.16), C_CYAN, alpha=55)
    deco_glow_dot(s, Inches(11.8), Inches(5.5), Inches(0.14), C_ORANGE, alpha=50)
    deco_glow_dot(s, Inches(1.5), Inches(6.5), Inches(0.12), C_VIOLET, alpha=45)
    # 大字 THANK YOU
    add_text(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.2),
             "THANK  YOU", size=80, color=C_TEXT_LT, bold=True, align=PP_ALIGN.CENTER)
    # 中文
    add_text(s, Inches(1.0), Inches(3.85), Inches(11.3), Inches(0.6),
             "感谢聆听  ·  敬请指导", size=30, color=C_CYAN, bold=True, align=PP_ALIGN.CENTER)
    # 分隔装饰线（双色）
    add_rect(s, Inches(5.4), Inches(4.75), Inches(1.2), Pt(2), color=C_ORANGE)
    add_rect(s, Inches(6.6), Inches(4.75), Inches(0.6), Pt(2), color=C_CYAN)
    # 项目名
    add_text(s, Inches(1.0), Inches(5.15), Inches(11.3), Inches(0.5),
             "AI 数字人工匠导师智能体  ·  职业教育全栈元宇宙解决方案",
             size=16, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)
    # 底部赛事
    add_text(s, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.4),
             "昌吉州 AI 智能体创新应用大赛  ·  2026",
             size=13, color=RGBColor(0x6A, 0x7A, 0x90), align=PP_ALIGN.CENTER)
    # 底部细线
    add_rect(s, Inches(4.5), Inches(6.7), Inches(4.3), Pt(0.5), color=C_LINE_DARK)


# ============ 组装 ============
slide_cover()
slide_toc()
slide_pain_points()
slide_goals()
slide_positioning()
slide_architecture()
slide_engine1()
slide_engine2()
slide_engine3()
slide_ai_agent()
slide_features()
slide_scenarios()
slide_pain_solution()
slide_quality()
slide_benefit()
slide_innovation()
slide_promotion()
slide_roadmap()
slide_ending()

out = r"d:\Admin_Platform_Project\admin-platform\市场资料\昌吉州AI智能体大赛-路演PPT.pptx"
prs.save(out)
print("PPT 已生成:", out)
print("总页数:", len(prs.slides._sldIdLst))
