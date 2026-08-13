# -*- coding: utf-8 -*-
"""将截图组装成PPT（每页为全屏图片，16:9）"""
from pptx import Presentation
from pptx.util import Inches
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

img_dir = r"d:\Admin_Platform_Project\admin-platform\市场资料\slides_img"
for i in range(1, 20):
    img = os.path.join(img_dir, f"slide_{i:02d}.png")
    slide = prs.slides.add_slide(BLANK)
    slide.shapes.add_picture(img, 0, 0, prs.slide_width, prs.slide_height)

out = r"d:\Admin_Platform_Project\admin-platform\市场资料\昌吉州AI智能体大赛-路演PPT-v3.pptx"
prs.save(out)
print("PPT 已生成:", out)
print("总页数:", len(prs.slides._sldIdLst))
